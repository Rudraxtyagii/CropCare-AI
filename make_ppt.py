from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "CropCare AI"
    subtitle.text = "Intelligent Plant Disease Detection System\n\nSummer Internship - 2026\nMajor Project"

    # 2. Problem Statement Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Problem Statement"
    tf = body_shape.text_frame
    tf.text = "Agriculture & Supply Chain Track"
    
    p = tf.add_paragraph()
    p.text = "Plant diseases cost the global economy billions annually and threaten food security."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Manual detection is slow, prone to error, and requires expert knowledge."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Solution: An automated, AI-driven computer vision system to detect diseases from leaf images instantly."
    p.level = 1

    # 3. Technical Architecture Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "System Architecture"
    tf = body_shape.text_frame
    tf.text = "End-to-End AI Application"
    
    p = tf.add_paragraph()
    p.text = "Frontend: Vanilla HTML/CSS/JS with Glassmorphism UI"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Backend: FastAPI (Python) for robust, high-speed API endpoints"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI/ML Model: PyTorch using MobileNetV2 (Transfer Learning)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Dataset: PlantVillage (38 classes of crop diseases)"
    p.level = 1

    # 4. Machine Learning Approach Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Machine Learning Approach"
    tf = body_shape.text_frame
    tf.text = "Model Details"
    
    p = tf.add_paragraph()
    p.text = "MobileNetV2: Chosen for optimal balance between accuracy and inference speed."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Preprocessing: Resizing (224x224), Center Cropping, ImageNet Normalization."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Loss Function: Cross-Entropy Loss for multi-class classification."
    p.level = 1

    # 5. Live Runtime Verification & QA Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Live Runtime Verification & QA"
    tf = body_shape.text_frame
    tf.text = "Empirical Proof of Functionality"
    
    p = tf.add_paragraph()
    p.text = "API Test Suite: Verified POST /predict endpoint live on HTTP port 8000."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "HTTP 200 OK: Valid JPEG/PNG leaf images return structured JSON with confidence tiers and treatments."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "HTTP 400 Bad Request: Successfully intercepts and rejects invalid text files and >5MB oversized images."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Dependency Audit: 100% clean package imports (FastAPI, PyTorch, Uvicorn, Pillow, Pydantic)."
    p.level = 1

    # 6. GitHub Repository & Deliverables Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "GitHub Repository & Deliverables"
    tf = body_shape.text_frame
    tf.text = "Open Source Project Repository"
    
    p = tf.add_paragraph()
    p.text = "GitHub URL: https://github.com/Rudraxtyagii/CropCare-AI"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI/ML Codebase: train.py, inference.py, and 38-class botanical knowledge base."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Production Server: main.py (FastAPI) with CORS and Pydantic validation."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Documentation: Technical Report, Verification QA Report, Title Page, and Summary PDF."
    p.level = 1

    # 7. Conclusion Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Conclusion & Future Work"
    tf = body_shape.text_frame
    tf.text = "Summary"
    
    p = tf.add_paragraph()
    p.text = "Successfully built an end-to-end AI application addressing a critical agricultural problem."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Future Work: Mobile App integration, real-time drone imagery analysis, and edge deployment."
    p.level = 1

    prs.save('Presentation.pptx')
    print("Presentation.pptx successfully created!")

if __name__ == '__main__':
    create_presentation()
